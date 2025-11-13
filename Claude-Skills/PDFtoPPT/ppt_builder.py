"""Utilities for assembling Monthly Business Review PowerPoints."""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, Iterable, List, Sequence, Tuple

from PIL import Image
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches

from pdf_extractor import SlideVisuals

LOGGER = logging.getLogger(__name__)

SLIDE_ORDER: Sequence[str] = [
    "(1) Summary",
    "(2) Total Co",
    "(3) Mkt. A",
    "(4) Mkt. B",
    "(5) Mkt. C",
]

EMUS_PER_INCH = 914400
MARGIN_IN = 0.6
GUTTER_IN = 0.4


def build_mbr_ppt(
    slide_visuals: Dict[str, SlideVisuals],
    template_path: Path,
    output_path: Path,
) -> None:
    """Populate the Monthly Business Review slides with extracted visuals."""

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    presentation = Presentation(template_path)
    slide_width_in = presentation.slide_width / EMUS_PER_INCH
    slide_height_in = presentation.slide_height / EMUS_PER_INCH

    for index, label in enumerate(SLIDE_ORDER):
        if index >= len(presentation.slides):
            LOGGER.warning(
                "Template only has %d slides; unable to fill slide %s",
                len(presentation.slides),
                label,
            )
            break
        slide = presentation.slides[index]
        _set_slide_title(slide, label)
        visuals = slide_visuals.get(label)
        if not visuals or not visuals.image_paths:
            LOGGER.info("No visuals available for %s; leaving slide mostly empty", label)
            continue
        if label == "(1) Summary":
            layout_boxes = _summary_layout(len(visuals.image_paths), slide_width_in, slide_height_in)
        else:
            layout_boxes = _dual_layout(len(visuals.image_paths), slide_width_in, slide_height_in)
        _place_images(slide, visuals.image_paths, layout_boxes)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    LOGGER.info("Saved PowerPoint to %s", output_path)


def _set_slide_title(slide: Slide, label: str) -> None:
    title_shape = slide.shapes.title
    if not title_shape or not title_shape.text_frame:
        return
    title_shape.text = _clean_title(label)


def _clean_title(label: str) -> str:
    if ")" in label:
        return label.split(")", 1)[1].strip()
    return label


def _summary_layout(count: int, slide_width: float, slide_height: float) -> List[Tuple[float, float, float, float]]:
    """Layout for the Summary slide (up to 3 visuals)."""

    usable_width = slide_width - 2 * MARGIN_IN
    usable_height = slide_height - 2 * MARGIN_IN
    boxes: List[Tuple[float, float, float, float]] = []

    if count <= 1:
        boxes.append((MARGIN_IN, MARGIN_IN, usable_width, usable_height))
        return boxes

    half_width = (usable_width - GUTTER_IN) / 2
    half_height = (usable_height - GUTTER_IN) / 2

    boxes.append((MARGIN_IN, MARGIN_IN, half_width, half_height))
    boxes.append((MARGIN_IN + half_width + GUTTER_IN, MARGIN_IN, half_width, half_height))

    if count >= 3:
        boxes.append((MARGIN_IN, MARGIN_IN + half_height + GUTTER_IN, usable_width, half_height))
    return boxes


def _dual_layout(count: int, slide_width: float, slide_height: float) -> List[Tuple[float, float, float, float]]:
    usable_width = slide_width - 2 * MARGIN_IN
    usable_height = slide_height - 2 * MARGIN_IN

    if count <= 1:
        return [(MARGIN_IN, MARGIN_IN, usable_width, usable_height)]

    half_width = (usable_width - GUTTER_IN) / 2
    return [
        (MARGIN_IN, MARGIN_IN, half_width, usable_height),
        (MARGIN_IN + half_width + GUTTER_IN, MARGIN_IN, half_width, usable_height),
    ]


def _place_images(
    slide: Slide,
    image_paths: Iterable[Path],
    layout_boxes: Sequence[Tuple[float, float, float, float]],
) -> None:
    for idx, image_path in enumerate(image_paths):
        if idx >= len(layout_boxes):
            LOGGER.warning("Too many visuals for slide; ignoring extra images")
            break
        left_in, top_in, width_in, height_in = layout_boxes[idx]
        _add_image(slide, image_path, left_in, top_in, width_in, height_in)


def _add_image(
    slide: Slide,
    image_path: Path,
    left_in: float,
    top_in: float,
    max_width_in: float,
    max_height_in: float,
) -> None:
    with Image.open(image_path) as img:
        img_width, img_height = img.size
    if img_width == 0 or img_height == 0:
        return
    width_ratio = max_width_in
    height_ratio = max_height_in
    target_width = min(width_ratio, height_ratio * (img_width / img_height))
    target_height = target_width * (img_height / img_width)
    if target_height > max_height_in:
        target_height = max_height_in
        target_width = target_height * (img_width / img_height)

    left = Inches(left_in + (max_width_in - target_width) / 2)
    top = Inches(top_in + (max_height_in - target_height) / 2)
    slide.shapes.add_picture(
        str(image_path),
        left=left,
        top=top,
        width=Inches(target_width),
    )


__all__ = ["build_mbr_ppt"]
